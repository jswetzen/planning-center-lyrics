#!/usr/bin/env python3
"""
pptx_export.py

Turns a Planning Center plan's songs into a PowerPoint deck for projection.

Like scheduler.py and live_session.py this has **no Flask in it** -- it takes
SongLyrics objects and returns bytes, so every decision here is unit-testable
without a server or network (tests/test_pptx_export.py). admin_app.py owns the
route and the plan lookup.

## Why .pptx and not .key

Keynote imports .pptx cleanly; there is no writable Keynote format. `.key` is
an undocumented, macOS-only bundle with no Python writer worth depending on,
so "export to Keynote" is served by handing Keynote a .pptx -- File > Open,
and it converts on the way in.

## One slide per stanza, not per song

The unit of projection is a screenful, not a song: nobody projects five
verses at once. Planning Center stores lyrics as plain text with
blank-line-separated stanzas, which is exactly that unit, so the split is the
document's own structure rather than something guessed.

Section labels ("Verse 1", "Refräng") are stripped, since those are notes for
the band and must not end up on a screen the congregation is reading. The
splitting, label detection and line wrapping all live in pco_client, shared
with the live projector so both surfaces agree on what a screenful is.

Repeated stanzas are deliberately **not** collapsed here, unlike on the
projector: a deck is advanced slide by slide, so a chorus sung three times
needs three slides.
"""

from __future__ import annotations

import io
import logging
import re
from typing import Iterable, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from pco_client import SongLyrics, split_stanzas, wrap_lines

log = logging.getLogger("pptx_export")

# 16:9. python-pptx still defaults to 4:3, which would letterbox every
# projector installed this century.
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

MARGIN = Inches(0.6)
FOOTER_HEIGHT = Inches(0.45)

# Font-size bounds for the lyric body. The floor matters more than the
# ceiling: below ~18pt a back-row reader is squinting, so an over-long stanza
# is better left slightly overflowing (and visibly needing a manual split)
# than silently shrunk to unreadable.
MIN_FONT_PT = 18
MAX_FONT_PT = 54
FOOTER_FONT_PT = 11

THEMES = {
    # Matches the projector's own default: white on black for a darkened room.
    "dark": {"bg": RGBColor(0x00, 0x00, 0x00), "fg": RGBColor(0xFF, 0xFF, 0xFF), "muted": RGBColor(0x77, 0x77, 0x77)},
    "light": {"bg": RGBColor(0xFF, 0xFF, 0xFF), "fg": RGBColor(0x00, 0x00, 0x00), "muted": RGBColor(0x88, 0x88, 0x88)},
}
DEFAULT_THEME = "dark"


def choose_font_size(stanza: str) -> int:
    """Pick a body font size that should fit the stanza on one slide.

    A deliberate estimate, not real text measurement: PowerPoint does its own
    layout with fonts this code can't see, and python-pptx's fit_text() needs
    font files present at render time (unreliable in a container). Estimating
    from the longest line and the line count gets close enough that decks come
    out consistent, and the caller can still nudge sizes in PowerPoint.

    Both constraints are applied and the smaller wins: width (a long line must
    not run off the side) and height (many lines must not run off the bottom).
    Callers pass text already through pco_client.wrap_lines, so the longest
    line here is one someone chose rather than whatever the source happened to
    contain.
    """
    lines = [line for line in stanza.split("\n") if line.strip()]
    if not lines:
        return MAX_FONT_PT

    usable_width_pt = (SLIDE_WIDTH - 2 * MARGIN) / Emu(1) * 72 / 914400
    usable_height_pt = (SLIDE_HEIGHT - 2 * MARGIN - FOOTER_HEIGHT) / Emu(1) * 72 / 914400

    longest = max(len(line) for line in lines)
    # ~0.5em per character averaged over a sans-serif face.
    by_width = usable_width_pt / (0.5 * longest) if longest else MAX_FONT_PT
    # 1.25 line spacing.
    by_height = usable_height_pt / (1.25 * len(lines))

    return int(max(MIN_FONT_PT, min(MAX_FONT_PT, by_width, by_height)))


def footer_text(song: SongLyrics) -> str:
    """The small credit line at the bottom of each slide.

    The CCLI number is on every slide on purpose: reporting what was projected
    is the licence-holder's obligation, and a number that only exists in
    Planning Center is one nobody will transcribe afterwards.
    """
    if song.ccli_number:
        return f"{song.title} · CCLI {song.ccli_number}"
    return song.title


# --------------------------------------------------------------------------
# Deck building
# --------------------------------------------------------------------------


def _add_slide(prs: Presentation, theme: dict, body: str, footer: str) -> None:
    blank_layout = prs.slide_layouts[6]  # 6 is "Blank" in the default template
    slide = prs.slides.add_slide(blank_layout)

    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = theme["bg"]

    box = slide.shapes.add_textbox(
        MARGIN, MARGIN, SLIDE_WIDTH - 2 * MARGIN, SLIDE_HEIGHT - 2 * MARGIN - FOOTER_HEIGHT
    )
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    size = choose_font_size(body)
    for index, line in enumerate(body.split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(size)
            run.font.bold = True
            run.font.color.rgb = theme["fg"]

    if footer:
        footer_box = slide.shapes.add_textbox(
            MARGIN, SLIDE_HEIGHT - MARGIN - FOOTER_HEIGHT, SLIDE_WIDTH - 2 * MARGIN, FOOTER_HEIGHT
        )
        footer_frame = footer_box.text_frame
        footer_frame.word_wrap = True
        paragraph = footer_frame.paragraphs[0]
        paragraph.text = footer
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = Pt(FOOTER_FONT_PT)
            run.font.color.rgb = theme["muted"]


def build_deck(
    songs: Iterable[SongLyrics], theme: str = DEFAULT_THEME, strip_labels: bool = True
) -> Presentation:
    """Build a 16:9 deck: one slide per stanza, in plan order."""
    palette = THEMES.get(theme, THEMES[DEFAULT_THEME])

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for song in songs:
        # Checked against the field rather than body(), which substitutes a
        # prose "_No lyrics found..._" note that reads fine in a Markdown
        # export and would be absurd projected at a congregation.
        stanzas = (
            [wrap_lines(st) for st in split_stanzas(song.plain_lyrics, strip_labels=strip_labels)]
            if song.plain_lyrics.strip()
            else []
        )
        if not stanzas:
            # A song with no lyrics in Planning Center still gets a slide, so
            # the running order stays intact and the gap is obvious in the
            # deck rather than being discovered mid-service.
            log.warning("No lyrics for %r; adding a placeholder slide.", song.title)
            _add_slide(prs, palette, f"[{song.title}]", footer_text(song))
            continue
        for stanza in stanzas:
            _add_slide(prs, palette, stanza, footer_text(song))

    return prs


def build_deck_bytes(
    songs: Iterable[SongLyrics], theme: str = DEFAULT_THEME, strip_labels: bool = True
) -> bytes:
    """build_deck, serialized -- what the download route actually sends."""
    buffer = io.BytesIO()
    build_deck(songs, theme=theme, strip_labels=strip_labels).save(buffer)
    return buffer.getvalue()


def suggest_filename(title_prefix: str, plan_date: Optional[str]) -> str:
    """`Lovsång Brokyrkan - 2026-07-26.pptx`, matching the Notion export's naming."""
    stem = f"{title_prefix} - {plan_date}" if plan_date else title_prefix
    # Windows and macOS both reject these outright in filenames.
    safe = re.sub(r'[<>:"/\\|?*]', "-", stem).strip()
    # A prefix made entirely of separators substitutes to "---", which is
    # non-empty but no more useful as a filename than "" was.
    if not safe.strip("-_ "):
        safe = "lyrics"
    return f"{safe}.pptx"
