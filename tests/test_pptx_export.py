"""
Tests for pptx_export.py -- the deck itself.

Stanza splitting, label stripping and line wrapping now live in pco_client and
are tested in tests/test_lyric_shaping.py; what's left here is what's specific
to PowerPoint: font sizing, slide structure, and the ways a deck differs from
the projector.
"""

import io

import pytest
from pptx import Presentation

import pptx_export
from pptx_export import (
    MAX_FONT_PT,
    MIN_FONT_PT,
    build_deck,
    build_deck_bytes,
    choose_font_size,
    footer_text,
    suggest_filename,
)
from pco_client import DEFAULT_WRAP_CHARS, SongLyrics


def _song(title="Song A", lyrics="line one\nline two", ccli="1234"):
    return SongLyrics(title, lyrics, "chords here", ccli, item_id="i1")


def _slide_texts(prs):
    """Every slide's body text, footers excluded."""
    out = []
    for slide in prs.slides:
        boxes = [sh for sh in slide.shapes if sh.has_text_frame]
        out.append(boxes[0].text_frame.text if boxes else "")
    return out


# --------------------------------------------------------------------------
# Font sizing
# --------------------------------------------------------------------------


def test_short_stanza_gets_the_largest_size():
    assert choose_font_size("Hey") == MAX_FONT_PT


def test_more_lines_never_increases_the_font_size():
    sizes = [choose_font_size("\n".join(["a line of lyrics"] * n)) for n in range(1, 15)]
    assert sizes == sorted(sizes, reverse=True)


def test_a_very_long_line_shrinks_the_font():
    assert choose_font_size("x" * 200) < choose_font_size("x" * 20)


def test_font_size_stays_within_bounds():
    for stanza in ["a", "\n".join(["long line of words here"] * 50), "y" * 1000]:
        assert MIN_FONT_PT <= choose_font_size(stanza) <= MAX_FONT_PT


def test_empty_stanza_does_not_divide_by_zero():
    assert choose_font_size("") == MAX_FONT_PT


# --------------------------------------------------------------------------
# Deck building
# --------------------------------------------------------------------------


def test_one_slide_per_stanza_in_plan_order():
    songs = [_song("A", "a1\n\na2"), _song("B", "b1")]
    assert _slide_texts(build_deck(songs)) == ["a1", "a2", "b1"]


def test_deck_is_widescreen():
    """python-pptx defaults to 4:3, which would letterbox on every projector
    installed this century."""
    prs = build_deck([_song()])
    assert round(prs.slide_width / prs.slide_height, 2) == round(16 / 9, 2)


def test_ccli_number_appears_on_every_slide():
    """Reporting what was projected is the licence-holder's obligation, and a
    number that only lives in Planning Center won't be transcribed later."""
    prs = build_deck([_song("A", "one\n\ntwo", ccli="7207484")])
    for slide in prs.slides:
        assert any("CCLI 7207484" in sh.text_frame.text for sh in slide.shapes if sh.has_text_frame)


def test_footer_omits_ccli_when_the_song_has_none():
    assert footer_text(_song("A", ccli=None)) == "A"


def test_song_without_lyrics_still_gets_a_placeholder_slide():
    """The gap should be obvious in the deck, not discovered mid-service."""
    texts = _slide_texts(build_deck([SongLyrics("Silent One", "", "", None, item_id="i9")]))
    assert texts == ["[Silent One]"]


def test_chords_are_never_exported():
    """A chord chart on the projector is the band's sheet on the wall."""
    song = SongLyrics("A", "just the words", "G  C  D\nchords", "1", item_id="i1")
    assert "chords" not in " ".join(_slide_texts(build_deck([song])))


def test_themes_differ_in_background():
    dark = build_deck([_song()], theme="dark")
    light = build_deck([_song()], theme="light")
    assert dark.slides[0].background.fill.fore_color.rgb != light.slides[0].background.fill.fore_color.rgb


def test_unknown_theme_falls_back_to_the_default():
    assert (
        build_deck([_song()], theme="chartreuse").slides[0].background.fill.fore_color.rgb
        == build_deck([_song()]).slides[0].background.fill.fore_color.rgb
    )


def test_empty_plan_produces_an_empty_but_valid_deck():
    prs = Presentation(io.BytesIO(build_deck_bytes([])))
    assert len(prs.slides) == 0


def test_bytes_round_trip_as_a_real_pptx():
    """The bytes the route hands the browser must actually open."""
    data = build_deck_bytes([_song("A", "one\n\ntwo")])
    assert data[:2] == b"PK"  # OOXML is a zip
    assert len(Presentation(io.BytesIO(data)).slides) == 2


# --------------------------------------------------------------------------
# Filenames
# --------------------------------------------------------------------------


def test_filename_matches_the_notion_export_naming():
    assert suggest_filename("Lovsång Brokyrkan", "2026-07-26") == "Lovsång Brokyrkan - 2026-07-26.pptx"


def test_filename_without_a_date():
    assert suggest_filename("Lovsång Brokyrkan", None) == "Lovsång Brokyrkan.pptx"


def test_filename_strips_characters_windows_and_macos_reject():
    assert "/" not in suggest_filename("A/B:C", "2026-01-01")
    assert suggest_filename('a<b>c:d"e/f\\g|h?i*j', None) == "a-b-c-d-e-f-g-h-i-j.pptx"


def test_filename_never_ends_up_empty():
    assert suggest_filename("///", None) == "lyrics.pptx"


# --------------------------------------------------------------------------
# Where a deck deliberately differs from the projector
# --------------------------------------------------------------------------


def test_repeated_stanzas_each_get_their_own_slide():
    """The projector collapses repeats to buy font size; a deck must not.
    It's advanced slide by slide, so a chorus sung three times needs three
    slides or the running order breaks mid-service."""
    song = _song("A", "sjung nu ut\n\nen ny sång\n\nsjung nu ut")
    assert _slide_texts(build_deck([song])) == ["sjung nu ut", "en ny sång", "sjung nu ut"]


def test_long_lines_are_wrapped_before_sizing():
    """Otherwise one runaway line drives the whole slide to the font floor."""
    long_line = "Du som är förlåten och du som blivit fri, sjung nu ut all ära till Guds lamm."
    text = _slide_texts(build_deck([_song("A", long_line)]))[0]
    assert "\n" in text
    assert all(len(l) <= DEFAULT_WRAP_CHARS for l in text.split("\n"))


def test_wrapping_raises_the_font_size_it_would_otherwise_force():
    long_line = "Du som är förlåten och du som blivit fri, sjung nu ut all ära till Guds lamm."
    assert choose_font_size(long_line) < choose_font_size(
        _slide_texts(build_deck([_song("A", long_line)]))[0]
    )


def test_section_labels_do_not_reach_a_slide():
    texts = _slide_texts(build_deck([_song("A", "VERSE 1:\nsing this\n\nCHORUS:\nand this")]))
    assert texts == ["sing this", "and this"]


def test_midblock_label_splits_into_two_slides():
    """Previously these merged into one slide with the marker printed in the
    middle of it."""
    assert _slide_texts(build_deck([_song("A", "line A\nINTRO\nline B")])) == ["line A", "line B"]
